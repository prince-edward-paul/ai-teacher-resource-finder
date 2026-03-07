import streamlit as st
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from pptx.util import Pt
import re, os, random, io, hashlib, time
from PIL import Image, ImageDraw, ImageFont
import pandas as pd
from datetime import datetime
from dotenv import load_dotenv

# -------------------------
# Load environment variables
# -------------------------
load_dotenv()
gemini_key = os.environ.get("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY")

# -------------------------
# Page setup
# -------------------------
st.set_page_config(page_title="AI Teacher Resource Finder", page_icon="🎓", layout="wide")
if "requests_today" not in st.session_state: st.session_state.requests_today = 0
if "last_request" not in st.session_state: st.session_state.last_request = 0
if "logged_in_teacher" not in st.session_state: st.session_state.logged_in_teacher = None

# -------------------------
# Folders & CSV DB setup
# -------------------------
TEMPLATE_DIR = "templates"
PREVIEW_DIR = "template_previews"
DOWNLOAD_LOG = "downloads.csv"
USERS_DB = "teachers.csv"
SAVED_RESOURCES_DB = "saved_resources.csv"
MARKETPLACE_DB = "marketplace.csv"
MARKETPLACE_FEEDBACK_DB = "marketplace_feedback.csv"

os.makedirs(TEMPLATE_DIR, exist_ok=True)
os.makedirs(PREVIEW_DIR, exist_ok=True)
available_templates = [f for f in os.listdir(TEMPLATE_DIR) if f.endswith(".pptx")]

# Ensure CSVs exist
for db_path, columns in [
    (USERS_DB, ["username","password_hash","email","role"]),
    (SAVED_RESOURCES_DB, ["teacher_username","resource_name","resource_type","file_path","timestamp"]),
    (MARKETPLACE_DB, ["teacher_username","resource_name","resource_type","file_path","timestamp"]),
    (MARKETPLACE_FEEDBACK_DB, ["resource_name","teacher_username","reviewer_username","rating","comment","timestamp"])
]:
    if not os.path.exists(db_path): pd.DataFrame(columns=columns).to_csv(db_path, index=False)

# -------------------------
# Helper functions
# -------------------------
def sanitize_filename(name): return re.sub(r"[^\w\s-]","",name).replace(" ","_")
def create_card_preview(text,color="#4B8BBE"):
    img = Image.new("RGB",(400,250),color); draw = ImageDraw.Draw(img)
    try: font = ImageFont.truetype("arial.ttf",28)
    except: font = ImageFont.load_default()
    text = text.replace("_"," ").title()
    bbox = draw.textbbox((0,0),text,font=font)
    w = bbox[2]-bbox[0]; h = bbox[3]-bbox[1]
    draw.text(((400-w)/2,(250-h)/2),text,fill="white",font=font)
    return img

def log_download(resource_name,user="anonymous"):
    row = pd.DataFrame([[resource_name,user,datetime.now()]],columns=["Resource","User","Timestamp"])
    if os.path.exists(DOWNLOAD_LOG): row.to_csv(DOWNLOAD_LOG,mode="a",header=False,index=False)
    else: row.to_csv(DOWNLOAD_LOG,index=False)

def save_resource_for_teacher(username,resource_name,resource_type,file_path):
    row = pd.DataFrame([[username,resource_name,resource_type,file_path,datetime.now()]],
        columns=["teacher_username","resource_name","resource_type","file_path","timestamp"])
    if os.path.exists(SAVED_RESOURCES_DB): row.to_csv(SAVED_RESOURCES_DB,mode="a",header=False,index=False)
    else: row.to_csv(SAVED_RESOURCES_DB,index=False)

def share_resource_to_marketplace(username,resource_name,resource_type,file_path):
    df_market = pd.read_csv(MARKETPLACE_DB)
    if ((df_market["teacher_username"]==username)&(df_market["resource_name"]==resource_name)).any():
        return False,"Resource already shared"
    row = pd.DataFrame([[username,resource_name,resource_type,file_path,datetime.now()]],
        columns=["teacher_username","resource_name","resource_type","file_path","timestamp"])
    df_market = pd.concat([df_market,row],ignore_index=True)
    df_market.to_csv(MARKETPLACE_DB,index=False)
    return True,"Resource shared to marketplace"

def add_feedback(resource_name,owner,reviewer,rating,comment):
    row = pd.DataFrame([[resource_name,owner,reviewer,rating,comment,datetime.now()]],
        columns=["resource_name","teacher_username","reviewer_username","rating","comment","timestamp"])
    df_fb = pd.read_csv(MARKETPLACE_FEEDBACK_DB)
    df_fb = pd.concat([df_fb,row],ignore_index=True)
    df_fb.to_csv(MARKETPLACE_FEEDBACK_DB,index=False)

def get_avg_rating(resource_name,owner):
    df_fb = pd.read_csv(MARKETPLACE_FEEDBACK_DB)
    df_r = df_fb[(df_fb["resource_name"]==resource_name)&(df_fb["teacher_username"]==owner)]
    if df_r.empty: return None
    return round(df_r["rating"].mean(),2)

# -------------------------
# Teacher Accounts
# -------------------------
def hash_password(password): return hashlib.sha256(password.encode()).hexdigest()
def register_teacher(username,password,email,role="teacher"):
    df = pd.read_csv(USERS_DB)
    if username in df["username"].values: return False,"Username exists"
    df = pd.concat([df,pd.DataFrame([[username,hash_password(password),email,role]],columns=df.columns)],ignore_index=True)
    df.to_csv(USERS_DB,index=False)
    return True,"Registered successfully"
def login_teacher(username,password):
    df = pd.read_csv(USERS_DB)
    user = df[(df["username"]==username)&(df["password_hash"]==hash_password(password))]
    if not user.empty: return True,user.iloc[0].to_dict()
    else: return False,"Invalid username/password"

# -------------------------
# Sidebar Navigation
# -------------------------
st.sidebar.title("Navigation")
menu = st.sidebar.radio("Menu",["Home","Download Analytics","Settings"])

# -------------------------
# Hero
# -------------------------
st.markdown("""
<div style="
background: linear-gradient(135deg,#4B8BBE,#306998);
padding:30px;
border-radius:15px;
color:white;
text-align:center;
margin-bottom:20px;">
<h1>🎓 AI Teacher Resource Finder</h1>
<p style="font-size:18px;">Create lesson plans and slides instantly using AI</p>
</div>
""",unsafe_allow_html=True)

# -------------------------
# Home Page
# -------------------------
if menu=="Home":
    if st.session_state.get("logged_in_teacher"):
        st.info(f"Logged in as: {st.session_state['logged_in_teacher']['username']}")
    else:
        st.warning("⚠️ Log in to track downloads & save resources.")

    st.subheader("🔍 Search Templates")
    search_query = st.text_input("Search template","")
    filtered_templates = [t for t in available_templates if search_query.lower() in t.lower()] if search_query else available_templates
    if search_query and not filtered_templates: st.info("No templates matched search")

    st.subheader("🎨 Choose a Template")
    cols = st.columns(3)
    for i, template_file in enumerate(filtered_templates+["🎲 Random Template"]):
        col = cols[i%3]
        base_name = os.path.splitext(template_file)[0]
        preview_path = os.path.join(PREVIEW_DIR,f"{base_name}.png")
        with col:
            if os.path.exists(preview_path): st.image(preview_path,use_container_width=True)
            else:
                img=create_card_preview(base_name)
                buf=io.BytesIO(); img.save(buf,format="PNG")
                st.image(buf.getvalue(),use_container_width=True)
            if st.button(f"Select {base_name}",key=f"btn_{i}"):
                st.session_state["chosen_template"]=template_file
                st.success(f"Selected {base_name}")

    chosen_template = st.session_state.get("chosen_template",None)
    topic = st.text_input("📘 Lesson Topic")

    if st.button("✨ Generate Lesson"):
        if not st.session_state.get("logged_in_teacher"):
            st.warning("Login first."); st.stop()
        if time.time()-st.session_state.last_request<10: st.warning("Wait 10s"); st.stop()
        st.session_state.last_request=time.time()
        if st.session_state.requests_today>=5: st.warning("Daily limit 5 reached"); st.stop()
        st.session_state.requests_today+=1
        if not topic: st.warning("Enter topic"); st.stop()
        if not chosen_template: st.warning("Select template"); st.stop()
        if not gemini_key: st.error("API key missing"); st.stop()

        with st.spinner("Generating lesson..."):
            prompt=f"""
Create a structured lesson plan for {topic}.
Include:
1. Objectives
2. Introduction
3. Key Teaching Points
4. Activities
5. Assessment
6. Summary
7. Slide Outline
"""
            for attempt in range(3):
                try:
                    genai.configure(api_key=gemini_key)
                    model=genai.GenerativeModel("gemini-2.0-flash")
                    lesson_text=model.generate_content(prompt).text
                    break
                except Exception as e:
                    if "ResourceExhausted" in str(e) and attempt<2: time.sleep(5)
                    elif "Quota exceeded" in str(e) or "429" in str(e):
                        st.error("Gemini API quota exceeded"); st.stop()
                    else: st.error(e); st.stop()

            safe_topic = sanitize_filename(topic)
            word_file=f"{safe_topic}_lesson.docx"; ppt_file=f"{safe_topic}_slides.pptx"
            doc=Document(); doc.add_heading(f"Lesson Plan: {topic}",level=1); doc.add_paragraph(lesson_text); doc.save(word_file)
            template_to_use=chosen_template if chosen_template!="🎲 Random Template" else random.choice(available_templates)
            prs=Presentation(os.path.join(TEMPLATE_DIR,template_to_use))
            sections=[s.strip() for s in lesson_text.split("\n\n") if s.strip()]
            for section in sections:
                slide=prs.slides.add_slide(prs.slide_layouts[1])
                lines=section.split("\n")
                slide.shapes.title.text=lines[0]
                if len(lines)>1: slide.placeholders[1].text="\n".join(lines[1:])
            prs.save(ppt_file)

            teacher_username = st.session_state["logged_in_teacher"]["username"]
            col1,col2=st.columns(2)
            with col1:
                with open(word_file,"rb") as f:
                    if st.download_button("📄 Download Word",f,file_name=word_file):
                        log_download(word_file,teacher_username)
                        save_resource_for_teacher(teacher_username,word_file,"Word",os.path.abspath(word_file))
            with col2:
                with open(ppt_file,"rb") as f:
                    if st.download_button("🎞️ Download Slides",f,file_name=ppt_file):
                        log_download(ppt_file,teacher_username)
                        save_resource_for_teacher(teacher_username,ppt_file,"PPT",os.path.abspath(ppt_file))

# -------------------------
# Analytics
# -------------------------
if menu=="Download Analytics":
    st.subheader("📊 Resource Downloads")
    if os.path.exists(DOWNLOAD_LOG):
        df=pd.read_csv(DOWNLOAD_LOG); st.dataframe(df); st.bar_chart(df["Resource"].value_counts())
    else: st.info("No downloads yet")

# -------------------------
# Settings & Marketplace + Teacher Dashboard
# -------------------------
if menu=="Settings":
    st.subheader("Teacher Accounts & Resources")
    if st.session_state["logged_in_teacher"]:
        teacher=st.session_state["logged_in_teacher"]
        st.success(f"✅ Logged in as {teacher['username']} ({teacher['role']})")
        if st.button("Logout"): st.session_state["logged_in_teacher"]=None; st.experimental_rerun()

        # ---------- Teacher Dashboard ----------
        st.markdown("### 📊 Teacher Dashboard")
        df_saved = pd.read_csv(SAVED_RESOURCES_DB)
        df_market = pd.read_csv(MARKETPLACE_DB)
        df_fb = pd.read_csv(MARKETPLACE_FEEDBACK_DB)
        user_saved = df_saved[df_saved["teacher_username"]==teacher["username"]]
        user_market = df_market[df_market["teacher_username"]==teacher["username"]]

        total_saved = len(user_saved)
        total_shared_downloads = df_saved.merge(pd.read_csv(DOWNLOAD_LOG), left_on="resource_name", right_on="Resource")
        total_shared_downloads = total_shared_downloads[total_shared_downloads["teacher_username"]==teacher["username"]].shape[0]
        avg_rating = df_fb[df_fb["teacher_username"]==teacher["username"]]["rating"].mean() if not df_fb[df_fb["teacher_username"]==teacher["username"]].empty else 0

        col1,col2,col3 = st.columns(3)
        col1.metric("Total Saved Resources",total_saved)
        col2.metric("Total Downloads of Shared Resources",total_shared_downloads)
        col3.metric("Average Rating of Shared Resources",round(avg_rating,2))

        # ---------- Downloads Over Time ----------
        st.markdown("#### 📈 Downloads Over Time")
        df_downloads = pd.read_csv(DOWNLOAD_LOG)
        teacher_downloads = df_downloads.merge(user_saved,left_on="Resource",right_on="resource_name")
        if not teacher_downloads.empty:
            downloads_time = teacher_downloads.groupby(teacher_downloads["Timestamp"].str[:10]).size()
            st.line_chart(downloads_time)
        else:
            st.info("No downloads yet.")

        # ---------- Rating Distribution ----------
        st.markdown("#### ⭐ Rating Distribution")
        teacher_ratings = df_fb[df_fb["teacher_username"]==teacher["username"]]
        if not teacher_ratings.empty:
            rating_dist = teacher_ratings["rating"].value_counts().sort_index()
            st.bar_chart(rating_dist)
        else:
            st.info("No ratings yet.")

        # ---------- Feedback List ----------
        st.markdown("#### 💬 Feedback on Your Resources")
        feedback = teacher_ratings
        if not feedback.empty:
            for _,row in feedback.iterrows():
                st.markdown(f"- {row['reviewer_username']} rated {row['resource_name']} {row['rating']}⭐ : {row['comment']}")
        else:
            st.info("No feedback yet.")

        # ---------- Tabs: Saved Resources + Marketplace ----------
        tabs=st.tabs(["💾 Saved Resources","🌐 Marketplace"])
        # Saved Resources
        with tabs[0]:
            st.markdown("### 💾 Your Saved Resources")
            teacher_resources = user_saved
            if not teacher_resources.empty:
                st.dataframe(teacher_resources[["resource_name","resource_type","timestamp"]])
                st.markdown("### Share to Marketplace")
                for _,row in teacher_resources.iterrows():
                    if st.button(f"Share {row['resource_name']}",key=f"share_{row['resource_name']}"):
                        success,msg=share_resource_to_marketplace(teacher["username"],row["resource_name"],row["resource_type"],row["file_path"])
                        st.success(msg) if success else st.warning(msg)
            else: st.info("No saved resources yet.")
        # Marketplace
        with tabs[1]:
            st.markdown("### 🌐 Community Marketplace")
            if not df_market.empty:
                resource_types=["All"] + sorted(df_market["resource_type"].unique().tolist())
                selected_type=st.selectbox("Filter by Resource Type",resource_types)
                sort_by=st.selectbox("Sort by",["Date","Rating"])
                search_keyword=st.text_input("Search resources by keyword","")

                df_disp=df_market.copy()
                if selected_type!="All": df_disp=df_disp[df_disp["resource_type"]==selected_type]
                if search_keyword: df_disp=df_disp[df_disp["resource_name"].str.contains(search_keyword,case=False)]

                df_disp["avg_rating"] = df_disp.apply(lambda x: get_avg_rating(x["resource_name"],x["teacher_username"]) or 0,axis=1)
                if sort_by=="Rating": df_disp=df_disp.sort_values("avg_rating",ascending=False)
                else: df_disp=df_disp.sort_values("timestamp",ascending=False)

                for _,row in df_disp.iterrows():
                    st.markdown(f"**{row['resource_name']}** by {row['teacher_username']} - Type: {row['resource_type']} - Avg Rating: {row['avg_rating'] if row['avg_rating']>0 else 'N/A'}")
                    with open(row['file_path'],"rb") as f: st.download_button("Download",f,file_name=row['resource_name'])
                    if teacher["username"]!=row["teacher_username"]:
                        rating=st.slider(f"Rate {row['resource_name']}",1,5,key=f"rate_{row['resource_name']}")
                        comment=st.text_input(f"Comment for {row['resource_name']}",key=f"comment_{row['resource_name']}")
                        if st.button(f"Submit Feedback",key=f"fb_{row['resource_name']}"):
                            add_feedback(row['resource_name'],row['teacher_username'],teacher["username"],rating,comment)
                            st.success("Feedback submitted!")
                    resource_fb=df_fb[(df_fb["resource_name"]==row['resource_name']) & (df_fb["teacher_username"]==row['teacher_username'])]
                    if not resource_fb.empty:
                        st.markdown("**Comments:**")
                        for _,c in resource_fb.iterrows():
                            st.markdown(f"- {c['reviewer_username']} ({c['rating']}⭐): {c['comment']}")
    else:
        tab1,tab2=st.tabs(["Login","Register"])
        with tab1:
            st.markdown("### Login")
            username=st.text_input("Username",key="login_username")
            password=st.text_input("Password",type="password",key="login_password")
            if st.button("Login",key="login_btn"):
                success,result=login_teacher(username,password)
                if success: st.session_state["logged_in_teacher"]=result; st.success(f"Welcome {username}"); st.experimental_rerun()
                else: st.error(result)
        with tab2:
            st.markdown("### Register")
            reg_username=st.text_input("Username",key="reg_username")
            reg_email=st.text_input("Email",key="reg_email")
            reg_password=st.text_input("Password",type="password",key="reg_password")
            if st.button("Register",key="reg_btn"):
                success,msg=register_teacher(reg_username,reg_password,reg_email)
                st.success(msg) if success else st.error(msg)

# -------------------------
# Footer
# -------------------------
st.sidebar.markdown("Developed by Prince Edward Paul © 2026")