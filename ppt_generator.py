# ppt_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import uuid
import re
import requests
from io import BytesIO

TEMPLATE_DIR = "templates"
GENERATED_DIR = "generated"
os.makedirs(GENERATED_DIR, exist_ok=True)

def sanitize_filename(name):
    return re.sub(r'[^a-zA-Z0-9_-]', '_', name)

def generate_presentation(topic, lesson_text, chosen_template=None):
    """
    Generates a student-centered PowerPoint presentation:
    - Uses chosen template or selects one randomly.
    - Inserts AI-generated lesson content slide by slide.
    - Adds suggested images from URLs.
    - 21st-century teaching standards (font, alignment, readability).
    Returns the path to the generated PPTX file.
    """

    # Select template
    template_file = chosen_template if chosen_template else \
        os.path.join(TEMPLATE_DIR, sorted(os.listdir(TEMPLATE_DIR))[0])
    prs = Presentation(os.path.join(TEMPLATE_DIR, template_file))

    # Split lesson into sections for slides
    sections = [s.strip() for s in lesson_text.split("\n\n") if s.strip()]

    for section in sections:
        lines = section.split("\n")
        slide_title = lines[0][:80]  # Limit title length for readability
        content_lines = lines[1:]  

        # Select slide layout (Title + Content)
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title

        # Insert content in placeholder or textbox
        content_text = "\n".join(content_lines)[:1000]  # Limit content size
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = content_text
            for para in slide.placeholders[1].text_frame.paragraphs:
                para.font.size = Pt(18)
                para.alignment = PP_ALIGN.LEFT
        else:
            left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = content_text
            for para in tf.paragraphs:
                para.font.size = Pt(18)
                para.alignment = PP_ALIGN.LEFT

        # Insert suggested images from content (e.g., "Image: https://...")
        for line in content_lines:
            if "Image:" in line or "Source:" in line:
                url_match = re.search(r"(https?://\S+)", line)
                if url_match:
                    url = url_match.group(1)
                    try:
                        response = requests.get(url)
                        img_stream = BytesIO(response.content)
                        # Add image to slide
                        slide.shapes.add_picture(img_stream, Inches(1), Inches(3), width=Inches(6))
                    except:
                        pass  # ignore image errors

    # Save presentation
    file_name = f"{sanitize_filename(topic)}_{uuid.uuid4().hex[:5]}.pptx"
    path = os.path.join(GENERATED_DIR, file_name)
    prs.save(path)
    return path